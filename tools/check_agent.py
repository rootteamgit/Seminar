#!/usr/bin/env python3
"""
セミナー再発防止エージェント
企画書（xlsx）のアジェンダと台本（docx）・スライド（pptx）の乖離をチェックし、
乖離があればSlackに通知する。

使い方:
  python check_agent.py <企画書.xlsx> <台本.docx> [スライド.pptx] [--回 N]

Claude Codeでの自動実行:
  git commit時にpre-commitフックから呼ばれる想定
"""

import sys
import os
import re
import json
import argparse
import textwrap
from pathlib import Path
from datetime import datetime

# ─────────────────────────────────────────────
# 1. テキスト抽出
# ─────────────────────────────────────────────

def extract_agenda_from_xlsx(xlsx_path: str, session_num: int = None) -> dict:
    """
    企画書xlsxからアジェンダを抽出する。
    セミナー企画書.xlsx の構造：
      - シート「企画書」またはシート名に回数を含む
      - 「アジェンダ」行を探して以降の項目を取得
      - 「テーマ」「登壇者」「ターゲット」も取得
    """
    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)

    result = {
        "theme": "",
        "speaker": "",
        "target": "",
        "agenda_items": [],
        "composition": "",   # 構成欄（台本生成の詳細指示）
        "lead_text": "",     # リード文
        "solve_issues": "",  # 解決できる課題
        "raw_text": []
    }

    # 対象シートを選択
    sheet = None
    for ws in wb.worksheets:
        name = ws.title
        # 回数指定がある場合は回数を含むシートを優先
        if session_num and (f"第{session_num}回" in name) and "TODO" not in name:
            sheet = ws
            break
    if sheet is None:
        # 「企画書」シートを探す
        for ws in wb.worksheets:
            if "企画書" in ws.title:
                sheet = ws
                break
    if sheet is None:
        sheet = wb.active

    for row in sheet.iter_rows(values_only=True):
        cells = [c for c in row if c is not None and str(c).strip() and str(c) != "None"]
        if not cells:
            continue
        row_text = " ".join(str(c) for c in cells)
        result["raw_text"].append(row_text)

        first = str(cells[0]).strip()
        second = str(cells[1]).strip() if len(cells) > 1 else ""

        # テーマ行
        if any(k in first for k in ["セミナーテーマ", "タイトル"]):
            result["theme"] = second

        # 登壇者行
        if "登壇者名" in first and "ふりがな" not in first:
            if result["speaker"] == "":
                result["speaker"] = second

        # ターゲット行
        if "ターゲット" in first:
            result["target"] = second

        # アジェンダ行：1セルに全アジェンダが改行で入っている
        if first == "アジェンダ" and second:
            for line in second.split("\n"):
                line = line.strip()
                if line and re.match(r'^[①②③④⑤⑥⑦⑧⑨⑩\d]', line):
                    result["agenda_items"].append(line)

        # 構成欄（常に取得・台本生成の詳細指示として使う）
        if "構成" in first and "ページ" in first and second:
            result["composition"] = second
        elif first == "構成" and second:
            result["composition"] = second
            # アジェンダが取れていない場合は構成欄から補完
            if not result["agenda_items"]:
                for line in second.split("\n"):
                    line = line.strip()
                    if line and re.match(r'^[①②③④⑤⑥⑦⑧⑨⑩\d]', line):
                        result["agenda_items"].append(line)

        # リード文
        if "リード文" in first and second:
            result["lead_text"] = second[:300]

        # 解決できる課題
        if "解決できる課題" in first and second:
            result["solve_issues"] = second[:300]

    wb.close()
    return result


def extract_text_from_docx(docx_path: str) -> str:
    """台本docx（またはmdテキスト）から全テキストを抽出"""
    import zipfile
    # ZIPかどうかで本物のdocxか判定
    try:
        with zipfile.ZipFile(docx_path):
            pass
        # 本物のdocx
        import docx as docx_lib
        doc = docx_lib.Document(docx_path)
        lines = []
        for para in doc.paragraphs:
            if para.text.strip():
                lines.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        lines.append(cell.text.strip())
        return "\n".join(lines)
    except (zipfile.BadZipFile, Exception):
        # テキスト/markdownファイルとして読み込み
        with open(docx_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()


def extract_text_from_pptx(pptx_path: str) -> list[str]:
    """スライドpptxから各スライドのタイトル・本文を抽出"""
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation(pptx_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides_text.append(f"[スライド{i}] " + " / ".join(texts[:5]))
    return slides_text

def extract_slide_refs_from_script(script_text: str) -> list[dict]:
    """台本から📊スライドN：タイトル の行を抽出して返す"""
    import re
    refs = []
    for line in script_text.split("\n"):
        # 「スライドN：タイトル」または「スライドN:タイトル」にマッチ
        m = re.search(r"スライド\s*(\d+)[：::]\s*(.+)", line)
        if m:
            refs.append({
                "num": int(m.group(1)),
                "title": m.group(2).strip()
            })
    return refs


def check_script_slide_alignment(script_text: str, slide_texts: list[str]) -> list[dict]:
    """
    台本のスライド番号とスライドファイルの枚数・タイトルを照合する。
    戻り値：乖離リスト [{severity, description}, ...]
    """
    issues = []
    script_refs = extract_slide_refs_from_script(script_text)

    if not script_refs:
        return []  # 台本にスライド番号がない場合はスキップ

    slide_count = len(slide_texts)
    script_max = max(r["num"] for r in script_refs) if script_refs else 0

    # 枚数チェック
    if slide_texts and script_max != slide_count:
        issues.append({
            "severity": "高",
            "description": (
                f"台本のスライド番号の最大値（S{script_max}）と"
                f"スライドファイルの枚数（{slide_count}枚）が一致しない。"
                f"台本のスライド番号を修正するか、スライドファイルを更新してください。"
            )
        })

    # 重複番号チェック
    nums = [r["num"] for r in script_refs]
    duplicates = [n for n in set(nums) if nums.count(n) > 1]
    if duplicates:
        issues.append({
            "severity": "中",
            "description": f"台本にスライド番号の重複があります：S{sorted(duplicates)}"
        })

    # 番号の連続性チェック（飛び番）→ セリフなしスライドの可能性があるため[高]
    if script_refs:
        unique_nums = sorted(set(nums))
        expected = list(range(1, unique_nums[-1] + 1))
        missing = [n for n in expected if n not in unique_nums]
        if missing:
            issues.append({
                "severity": "高",
                "description": (
                    f"台本にスライド番号の飛び番があります：S{missing}。"
                    f"これらのスライドに対応するセリフが台本に存在しない可能性があります。"
                    f"各スライドの切り替えタイミングで話す内容が台本に含まれているか確認してください。"
                )
            })

    # 台本バージョン番号チェック（タイトル行のv番号が古くないか）
    import re
    version_match = re.search(r"台本（v(\d+)）", script_text)
    if version_match:
        # バージョン番号が取得できた場合、ファイル名と一致するかは呼び出し側で確認
        # ここでは記録のみ（将来的にファイル名と照合する拡張ポイント）
        pass

    return issues


def check_script_version_title(script_text: str, docx_path: str) -> list[dict]:
    """
    台本内部のタイトル（v番号）とファイル名のv番号が一致するかチェック。
    例：ファイル名が _v8.docx なのに内部タイトルが（v7）なら[中]で検出。
    """
    import re, os
    issues = []

    # ファイル名からバージョン番号を抽出
    filename = os.path.basename(docx_path)
    file_ver_m = re.search(r"_v(\d+)\.docx", filename)
    if not file_ver_m:
        return issues
    file_ver = int(file_ver_m.group(1))

    # 台本テキストから内部タイトルのバージョン番号を抽出
    title_ver_m = re.search(r"台本（v(\d+)）", script_text)
    if not title_ver_m:
        return issues
    title_ver = int(title_ver_m.group(1))

    if file_ver != title_ver:
        issues.append({
            "severity": "中",
            "description": (
                f"ファイル名のバージョン（v{file_ver}）と"
                f"台本内部タイトルのバージョン（v{title_ver}）が一致しない。"
                f"台本の冒頭タイトルを「台本（v{file_ver}）」に更新してください。"
            )
        })

    return issues


# ─────────────────────────────────────────────
# 2. Claude API で乖離チェック
# ─────────────────────────────────────────────

def check_deviation_with_claude(
    agenda_data: dict,
    script_text: str,
    slide_texts: list[str],
    session_label: str
) -> dict:
    """
    Claude API（claude-sonnet-4-6）を使って乖離を判定する。
    戻り値: {"ok": bool, "issues": [...], "summary": str}
    """
    import urllib.request

    agenda_str = "\n".join(f"  - {item}" for item in agenda_data["agenda_items"]) or "（取得できませんでした）"
    composition_str = agenda_data.get("composition", "")
    lead_str = agenda_data.get("lead_text", "")
    solve_str = agenda_data.get("solve_issues", "")
    slide_str = "\n".join(slide_texts[:30]) if slide_texts else "（スライドなし）"
    # 台本は長いので最初の3000文字
    script_excerpt = script_text[:3000] if script_text else "（台本なし）"

    # APIに渡す前にテキストをサニタイズ（特殊文字・制御文字を除去）
    def sanitize(text):
        if not text:
            return ""
        import unicodedata
        # 制御文字を除去（タブ・改行は保持）
        result = "".join(c for c in text if unicodedata.category(c) != "Cc" or c in "\t\n")
        # バックスラッシュをエスケープ
        result = result.replace("\\", "\\\\")
        return result[:2000]  # 長すぎる場合は切り詰め（APIトークン節約）

    agenda_str     = sanitize(agenda_str)
    composition_str = sanitize(composition_str)
    lead_str       = sanitize(lead_str)
    solve_str      = sanitize(solve_str)
    script_excerpt = sanitize(script_excerpt)
    slide_str      = sanitize(slide_str)

    prompt = f"""あなたはセミナー制作の品質チェック担当です。
以下の「企画書アジェンダ・構成」と「台本・スライドの内容」を照合し、乖離・不整合を検出してください。

## 対象セミナー
{session_label}
テーマ：{agenda_data['theme']}
登壇者：{agenda_data['speaker']}

## 企画書のアジェンダ（正とする）
{agenda_str}

## 企画書の構成欄（台本の詳細指示）
{composition_str if composition_str else "（記載なし）"}

## 企画書のリード文（ターゲットの課題感）
{lead_str if lead_str else "（記載なし）"}

## 企画書の解決できる課題
{solve_str if solve_str else "（記載なし）"}

## 台本の冒頭抜粋（最初3000文字）
{script_excerpt}

## スライドのテキスト抜粋（各スライドの要素）
{slide_str}

## 判定基準
1. 台本・スライドに、企画書のアジェンダにない章・トピックが含まれていないか？
2. 企画書のアジェンダ項目が台本・スライドで欠落していないか？
3. 企画書の構成欄の指示（登壇形式・各章の内容指示）が台本に反映されているか？
4. 企画書のリード文・解決できる課題が台本の内容にカバーされているか？
5. 登壇者・テーマが台本・スライドと一致しているか？

## 出力形式（JSONのみ・マークダウン不要）
{{
  "ok": true or false,
  "issues": [
    {{"severity": "高/中/低", "description": "問題の説明（具体的に）"}}
  ],
  "summary": "全体の判定サマリー（2〜3文）"
}}

乖離がなければ "ok": true、issues: [] としてください。"""

    # promptをjson.dumpsで安全にシリアライズ
    payload = json.dumps({
        "model": "claude-sonnet-4-6",
        "max_tokens": 2000,
        "messages": [{"role": "user", "content": prompt}]
    }, ensure_ascii=False).encode("utf-8")

    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    headers = {
        "Content-Type": "application/json",
        "anthropic-version": "2023-06-01"
    }
    if api_key:
        headers["x-api-key"] = api_key
    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data=payload,
        headers=headers
    )

    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            body = json.loads(resp.read().decode("utf-8"))
            text = body["content"][0]["text"]
            # JSON部分だけ取り出す（マークダウンコードブロックを除去）
            text = re.sub(r"```json\s*", "", text)
            text = re.sub(r"```\s*", "", text)
            text = text.strip()
            return json.loads(text)
    except Exception as e:
        return {
            "ok": False,
            "issues": [{"severity": "高", "description": f"Claude API呼び出しエラー: {e}"}],
            "summary": "チェックAPIへの接続に失敗しました。"
        }


# ─────────────────────────────────────────────
# 3. Slack通知
# ─────────────────────────────────────────────

# ─────────────────────────────────────────────
# 4. 自動再発防止ロジック
# ─────────────────────────────────────────────

def auto_fix_add_detection(issues: list[dict], agent_path: str) -> list[str]:
    """
    検出された問題のうち、まだcheck_agent.pyに検出ロジックがないものを
    Claude APIを使って自動生成し、check_agent.pyに追記する。
    戻り値：追加した検出ロジックの説明リスト
    """
    import urllib.request, json, os, re

    if not issues:
        return []

    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        return []

    agent_code = open(agent_path).read()
    added = []

    for issue in issues:
        # severity=高 のみ自動修正対象
        if issue.get("severity") != "高":
            continue

        desc = issue.get("description", "")

        # 既にcheck_agentに類似の検出ロジックがあるか確認
        prompt = f"""あなたはPythonエンジニアです。
以下のcheck_agent.pyに、下記の問題を検出するロジックがすでに含まれているか確認してください。

## 検出したい問題
{desc}

## check_agent.pyの現在のコード（抜粋）
{agent_code[:6000]}

## 回答形式（JSONのみ）
{{"already_exists": true/false, "reason": "理由を1行で"}}
"""

        try:
            payload = json.dumps({
                "model": "claude-sonnet-4-6",
                "max_tokens": 300,
                "messages": [{"role": "user", "content": prompt}]
            }, ensure_ascii=False).encode("utf-8")

            req = urllib.request.Request(
                "https://api.anthropic.com/v1/messages",
                data=payload,
                headers={
                    "Content-Type": "application/json",
                    "anthropic-version": "2023-06-01",
                    "x-api-key": api_key
                }
            )
            with urllib.request.urlopen(req, timeout=30) as resp:
                body = json.loads(resp.read().decode("utf-8"))
                text = body["content"][0]["text"]
                text = re.sub(r"```json\s*", "", text)
                text = re.sub(r"```\s*", "", text)
                check = json.loads(text.strip())

            if check.get("already_exists"):
                continue

            # 新しい検出ロジックをClaude APIで生成
            fix_prompt = f"""あなたはPythonエンジニアです。
以下の問題を検出するPython関数を1つ書いてください。

## 検出したい問題
{desc}

## 関数の仕様
- 関数名: check_auto_[問題を表す英語スネークケース]
- 引数: script_text: str, slide_texts: list[str], agenda_data: dict
- 戻り値: list[dict] （各dictは {{"severity": "高"/"中"/"低", "description": "問題の説明"}} の形式）
- 問題がなければ空リストを返す
- シンプルに、文字列マッチングや数値比較だけで実装する
- Claude API呼び出しは不要

Pythonコードのみ出力してください。説明不要。
```python
def check_auto_...
```
"""
            payload2 = json.dumps({
                "model": "claude-sonnet-4-6",
                "max_tokens": 800,
                "messages": [{"role": "user", "content": fix_prompt}]
            }, ensure_ascii=False).encode("utf-8")

            req2 = urllib.request.Request(
                "https://api.anthropic.com/v1/messages",
                data=payload2,
                headers={
                    "Content-Type": "application/json",
                    "anthropic-version": "2023-06-01",
                    "x-api-key": api_key
                }
            )
            with urllib.request.urlopen(req2, timeout=30) as resp2:
                body2 = json.loads(resp2.read().decode("utf-8"))
                new_code = body2["content"][0]["text"]
                # コードブロックを取り出す
                m = re.search(r"```python\n(.+?)```", new_code, re.DOTALL)
                if m:
                    new_code = m.group(1)

            # check_agent.pyの末尾（if __name__ の前）に追記
            insert_marker = "\nif __name__ == \"__main__\":"
            if insert_marker in agent_code and new_code.strip().startswith("def check_auto_"):
                updated = agent_code.replace(
                    insert_marker,
                    f"\n\n# [自動追加 再発防止ロジック]\n{new_code}\n{insert_marker}"
                )
                open(agent_path, 'w').write(updated)
                agent_code = updated  # 次のループで最新コードを参照
                func_name = re.search(r"def (check_auto_\w+)", new_code)
                fname = func_name.group(1) if func_name else "unknown"
                added.append(f"{fname}：{desc[:50]}")

        except Exception as e:
            print(f"[auto-fix] スキップ（{e}）")
            continue

    return added


def notify_slack(webhook_url: str, result: dict, session_label: str, files: dict):
    """乖離チェック結果をSlackに投稿する"""
    import urllib.request

    if result["ok"]:
        icon = "✅"
        color = "#36a64f"
        header = f"{icon} 企画書との乖離なし：{session_label}"
    else:
        high_count = sum(1 for i in result["issues"] if i["severity"] == "高")
        icon = "🚨" if high_count > 0 else "⚠️"
        color = "#ff0000" if high_count > 0 else "#ffa500"
        header = f"{icon} 乖離を検出：{session_label}（{len(result['issues'])}件）"

    # issues を整形
    issues_text = ""
    for issue in result.get("issues", []):
        sev = issue["severity"]
        sev_icon = {"高": "🔴", "中": "🟡", "低": "🔵"}.get(sev, "•")
        issues_text += f"{sev_icon} [{sev}] {issue['description']}\n"

    files_text = "\n".join(f"• {k}: `{v}`" for k, v in files.items())
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")

    blocks = [
        {
            "type": "header",
            "text": {"type": "plain_text", "text": header, "emoji": True}
        },
        {
            "type": "section",
            "text": {"type": "mrkdwn", "text": f"*サマリー：* {result['summary']}"}
        }
    ]

    if issues_text:
        blocks.append({
            "type": "section",
            "text": {"type": "mrkdwn", "text": f"*検出された問題：*\n{issues_text.strip()}"}
        })

    blocks += [
        {"type": "divider"},
        {
            "type": "section",
            "text": {"type": "mrkdwn", "text": f"*チェック対象ファイル：*\n{files_text}\n\n_実行日時: {timestamp}_"}
        }
    ]

    payload = json.dumps({"blocks": blocks}).encode("utf-8")
    req = urllib.request.Request(
        webhook_url,
        data=payload,
        headers={"Content-Type": "application/json"}
    )
    try:
        urllib.request.urlopen(req, timeout=10)
        print("[Slack] 通知送信完了")
    except Exception as e:
        print(f"[Slack] 通知失敗: {e}", file=sys.stderr)


# ─────────────────────────────────────────────
# 4. メイン
# ─────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="セミナー再発防止エージェント")
    parser.add_argument("xlsx",   help="企画書 .xlsx のパス")
    parser.add_argument("docx",   help="台本 .docx のパス")
    parser.add_argument("pptx",   nargs="?", default=None, help="スライド .pptx のパス（省略可）")
    parser.add_argument("--回",   type=int, default=None, dest="session_num", help="セミナー回数（例: 21）")
    parser.add_argument("--slack-webhook", default=os.environ.get("SLACK_WEBHOOK_URL", ""),
                        help="Slack Incoming Webhook URL（環境変数 SLACK_WEBHOOK_URL でも設定可）")
    parser.add_argument("--dry-run", action="store_true", help="Slack通知せず結果をターミナルに出力")
    args = parser.parse_args()

    # ─── pre-flight check ───
    # スライドファイルなしで台本生成→チェックしようとしていないか警告
    if not args.pptx:
        print("[pre-flight] ⚠️  スライドファイルが指定されていません。")
        print("             スライドを先に完成させてから台本チェックを行うことを推奨します。")
        print("             スライドなしで続行する場合は --skip-slide-check を追加してください。")
        # --skip-slide-check がない場合は確認を求める（CI環境では自動スキップ）
        import sys as _sys
        if '--skip-slide-check' not in _sys.argv and not os.environ.get("CI"):
            ans = input("             スライドなしで続行しますか？ [y/N]: ").strip().lower()
            if ans != 'y':
                print("[pre-flight] 中断しました。スライドファイルを用意してから再実行してください。")
                sys.exit(1)

    # ファイル存在確認
    for path, label in [(args.xlsx, "企画書"), (args.docx, "台本")]:
        if not Path(path).exists():
            print(f"[エラー] {label}ファイルが見つかりません: {path}", file=sys.stderr)
            sys.exit(1)

    session_label = f"第{args.session_num}回" if args.session_num else Path(args.docx).stem

    print(f"[チェック開始] {session_label}")
    print(f"  企画書: {args.xlsx}")
    print(f"  台本  : {args.docx}")
    if args.pptx:
        print(f"  スライド: {args.pptx}")

    # テキスト抽出
    print("[1/3] テキスト抽出中...")
    agenda_data = extract_agenda_from_xlsx(args.xlsx, args.session_num)
    script_text = extract_text_from_docx(args.docx)
    slide_texts = extract_text_from_pptx(args.pptx) if args.pptx and Path(args.pptx).exists() else []

    if not agenda_data["agenda_items"]:
        print("[警告] 企画書からアジェンダ項目を抽出できませんでした。raw_text で代替します。")
        # raw_text をフォールバック用アジェンダとして使用
        agenda_data["agenda_items"] = agenda_data["raw_text"][:20]

    print(f"  → アジェンダ項目: {len(agenda_data['agenda_items'])}件")
    print(f"  → 台本テキスト: {len(script_text)}文字")
    print(f"  → スライド枚数: {len(slide_texts)}枚")

    # 台本バージョン番号チェック（ファイル名と内部タイトルの一致）
    version_issues = check_script_version_title(script_text, args.docx)
    if version_issues:
        print(f"[バージョン番号チェック] ⚠️ {len(version_issues)}件の不一致を検出")
        for iss in version_issues:
            print(f"  [{iss['severity']}] {iss['description']}")
    else:
        print("[バージョン番号チェック] ✅ バージョン番号OK")

    # 台本↔スライド対応チェック（ローカル・高速）
    alignment_issues = check_script_slide_alignment(script_text, slide_texts)
    if alignment_issues:
        print(f"[台本↔スライド照合] ⚠️ {len(alignment_issues)}件の対応ズレを検出")
        for iss in alignment_issues:
            print(f"  [{iss['severity']}] {iss['description']}")
    else:
        if slide_texts:
            print("[台本↔スライド照合] ✅ スライド番号の対応OK")

    # Claude API でチェック
    print("[2/3] Claude API でチェック中...")
    result = check_deviation_with_claude(agenda_data, script_text, slide_texts, session_label)

    # alignment_issues・version_issuesをresultにマージ
    all_local_issues = version_issues + alignment_issues
    if all_local_issues:
        result["issues"] = all_local_issues + result.get("issues", [])
        result["ok"] = False
        result["summary"] = f"ローカルチェック{len(all_local_issues)}件。" + result.get("summary", "")

    # 結果表示
    print("\n─── チェック結果 ───")
    if result["ok"]:
        print("✅ 乖離なし")
    else:
        print(f"⚠️  {len(result['issues'])}件の問題を検出")
        for issue in result["issues"]:
            print(f"  [{issue['severity']}] {issue['description']}")
    print(f"サマリー: {result['summary']}")
    print("─────────────────────\n")

    # ─── 自動再発防止ロジック追加 ───
    if not result["ok"]:
        high_issues = [i for i in result["issues"] if i["severity"] == "高"]
        if high_issues:
            print("[auto-fix] 高severityの問題を検出。再発防止ロジックを自動追加します...")
            agent_path = os.path.abspath(__file__)
            added = auto_fix_add_detection(high_issues, agent_path)
            if added:
                print(f"[auto-fix] ✅ {len(added)}件の検出ロジックを追加しました：")
                for a in added:
                    print(f"  - {a}")
                # 更新したcheck_agent.pyをgit add
                import subprocess
                try:
                    subprocess.run(["git", "-C", str(Path(agent_path).parent.parent),
                                    "add", agent_path], check=False)
                    print("[auto-fix] git add 完了")
                except Exception:
                    pass
            else:
                print("[auto-fix] 既存ロジックで対応済み（追加不要）")

    # Slack通知
    files_info = {"企画書": args.xlsx, "台本": args.docx}
    if args.pptx:
        files_info["スライド"] = args.pptx

    if args.dry_run:
        print("[dry-run] Slack通知はスキップしました")
    elif args.slack_webhook:
        print("[3/3] Slack通知中...")
        notify_slack(args.slack_webhook, result, session_label, files_info)
    else:
        print("[3/3] SLACK_WEBHOOK_URL が未設定のためSlack通知をスキップ")
        print("      設定方法: export SLACK_WEBHOOK_URL=https://hooks.slack.com/services/...")

    # 乖離あり → exit code 1（pre-commitフックでcommitブロックに使う）
    if not result["ok"]:
        high_issues = [i for i in result["issues"] if i["severity"] == "高"]
        if high_issues:
            print("[pre-commit] 重大な乖離があるためcommitを中断します。", file=sys.stderr)
            sys.exit(1)

    sys.exit(0)


if __name__ == "__main__":
    main()
